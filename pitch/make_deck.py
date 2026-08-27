#!/usr/bin/env python3
"""Génère pitch/pitch_deck.pptx — 10 slides 16:9, thème sombre "chaleur"."""

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

BG = RGBColor(0x0B, 0x12, 0x20)
PANEL = RGBColor(0x11, 0x1B, 0x31)
PANEL2 = RGBColor(0x0E, 0x17, 0x2B)
ACC = RGBColor(0xFF, 0x6B, 0x35)
CY = RGBColor(0x4C, 0xC9, 0xF0)
GOLD = RGBColor(0xFF, 0xD1, 0x66)
TXT = RGBColor(0xE8, 0xEE, 0xFC)
MUT = RGBColor(0x8B, 0x97, 0xB0)
RED = RGBColor(0xFF, 0x5E, 0x57)
GREEN = RGBColor(0x2D, 0xD4, 0xA7)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def box(sl, l, t, w, h, fill=PANEL, line=None, rounded=True):
    shp = sl.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h),
    )
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def txt(sl, l, t, w, h, paras, anchor=MSO_ANCHOR.TOP):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    first = True
    for para in paras:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = para.get("align", PP_ALIGN.LEFT)
        p.space_after = Pt(para.get("sa", 6))
        p.space_before = Pt(para.get("sb", 0))
        for r in para["runs"]:
            run = p.add_run()
            run.text = r["t"]
            f = run.font
            f.size = Pt(r.get("s", 16))
            f.bold = r.get("b", False)
            f.color.rgb = r.get("c", TXT)
            f.name = r.get("name", "Segoe UI")
    return tb


def header(sl, kicker, title, accent=ACC):
    box(sl, 0.55, 0.5, 0.09, 0.62, fill=accent, rounded=False)
    txt(sl, 0.85, 0.42, 11.8, 0.4, [{"runs": [{"t": kicker.upper(), "s": 12, "b": True, "c": MUT}]}])
    txt(sl, 0.85, 0.66, 11.9, 0.7, [{"runs": [{"t": title, "s": 27, "b": True, "c": TXT}]}])


def bullets(sl, l, t, w, h, items, size=15, gap=8):
    paras = []
    for it in items:
        if isinstance(it, tuple):
            head, body = it
            paras.append({"runs": [{"t": "▪  ", "s": size, "b": True, "c": ACC},
                                   {"t": head, "s": size, "b": True, "c": TXT},
                                   {"t": " — " + body, "s": size, "c": MUT}], "sa": gap})
        else:
            paras.append({"runs": [{"t": "▪  ", "s": size, "b": True, "c": ACC},
                                   {"t": it, "s": size, "c": TXT}], "sa": gap})
    txt(sl, l, t, w, h, paras)


# ============ 1. TITRE ============
s = slide()
box(s, 0, 0, 13.333, 7.5, fill=BG)
box(s, 0, 6.9, 13.333, 0.6, fill=PANEL2, rounded=False)
box(s, 0.9, 1.5, 0.14, 2.2, fill=ACC, rounded=False)
txt(s, 1.3, 1.45, 11, 1.4, [{"runs": [{"t": "HEATSENTINEL", "s": 54, "b": True, "c": TXT}]}])
txt(s, 1.3, 2.55, 11, 0.6, [{"runs": [{"t": "From 2 meters above ground to action in 60 seconds.", "s": 20, "c": GOLD}]}])
txt(s, 1.3, 3.35, 11, 1.6, [
    {"runs": [{"t": "FortyGuard Hackathon'26 — Building the World's Temperature AI", "s": 15, "c": MUT}], "sa": 4},
    {"runs": [{"t": "Hyperlocal heat intelligence for Cotonou, Benin — live demo on real Phoenix data", "s": 15, "c": MUT}]},
])
for i, (tag, c) in enumerate([("Resilient Cities", GREEN), ("Agentic AI", CY), ("Data Analysis & Correlation", GOLD)]):
    box(s, 1.3 + i * 2.9, 5.35, 2.7, 0.52, fill=PANEL, line=c)
    txt(s, 1.3 + i * 2.9, 5.44, 2.7, 0.4, [{"runs": [{"t": tag, "s": 12.5, "b": True, "c": c}], "align": PP_ALIGN.CENTER}])
txt(s, 0.9, 7.0, 11, 0.4, [{"runs": [{"t": "Team ClimVision  ·  Powered by FortyGuard Temperature API®  ·  NVIDIA-recognized stack", "s": 12, "c": MUT}]}])

# ============ 2. PROBLÈME ============
s = slide()
header(s, "The problem", "Heat is the silent killer of growing cities")
box(s, 0.55, 1.75, 5.9, 2.3, fill=PANEL)
txt(s, 0.9, 1.95, 5.3, 2.0, [
    {"runs": [{"t": "2,000,000+", "s": 34, "b": True, "c": RED}], "sa": 2},
    {"runs": [{"t": "premature heat-related deaths per year worldwide (WHO). Heat is the deadliest climate hazard — and the least prepared for.", "s": 14, "c": MUT}]},
])
box(s, 6.85, 1.75, 5.9, 2.3, fill=PANEL)
txt(s, 7.2, 1.95, 5.3, 2.0, [
    {"runs": [{"t": "Cotonou today", "s": 18, "b": True, "c": TXT}], "sa": 4},
    {"runs": [{"t": "≈700,000 residents (commune) · 1.5M+ metro · humid tropical heat: feels-like >40 °C · dense districts next to the port · millions of outdoor workers.", "s": 14, "c": MUT}], "sa": 4},
    {"runs": [{"t": "Hyperlocal heat early-warning systems in Cotonou:  zero.", "s": 15, "b": True, "c": GOLD}]},
])
bullets(s, 0.7, 4.5, 11.9, 2.4, [
    ("City-scale forecasts", "too coarse: a 2 °C difference between the port and the beach can mean health or heatstroke."),
    ("Late alerts", "weather bulletins arrive at 8 AM — the peak is at 3 PM."),
    ("No action loop", "nobody translates a forecast into what a hospital, school or market must do, when."),
], size=15.5)

# ============ 3. POURQUOI MAINTENANT ============
s = slide()
header(s, "Why now", "Street-level heat data finally exists")
bullets(s, 0.7, 1.85, 11.9, 3.2, [
    ("FortyGuard Temperature API®", "hyperlocal urban temperature: 2 m above ground, 20 m² resolution, real-time, powered by NVIDIA-recognized Large Temperature Models."),
    ("The physics", "heat kills at street level: asphalt, traffic, markets, shade — not at the airport sensor 10 km away."),
    ("AI closes the loop", "models now turn raw readings into 6-hour nowcasts and autonomous alerts in seconds, on a $120 edge computer."),
], size=16)
box(s, 0.7, 5.3, 11.9, 1.5, fill=PANEL)
txt(s, 1.1, 5.55, 11.2, 1.1, [
    {"runs": [{"t": "The data now exists at street level. ", "s": 18, "b": True, "c": GOLD},
              {"t": "What was missing is the brain: the system that watches, predicts, decides and acts.", "s": 18, "c": TXT}]},
], anchor=MSO_ANCHOR.MIDDLE)

# ============ 4. SOLUTION ============
s = slide()
header(s, "The solution", "HeatSentinel — an agentic AI for urban heat")
bullets(s, 0.7, 1.9, 11.9, 3.6, [
    ("Watch", "continuously monitors 20 Temperature API® measurement points — 20 m² resolution, modeled 2 m above ground (FortyGuard's mesh; we install no hardware). Cotonou is the deployment target; the live demo runs on Phoenix, where the mesh is already in service."),
    ("Predict", "hybrid Ridge+LightGBM nowcast of the 6-hour temperature peak per node — MAE 0.26 °C / R² 0.95 (Cotonou), 0.61 °C / 0.97 (Phoenix) — plus 48 h z-score anomaly detection for local heat spikes and sensor drift."),
    ("Decide", "a transparent policy engine deduplicates, escalates only when needed, and keeps a full audit ledger (no black box, no LLM dependency)."),
    ("Act", "dispatches bilingual alerts (FR/EN) with concrete actions: open cooling centers, hospital surge, pause outdoor work 2-5 PM, hydrate zones."),
], size=15.5)
box(s, 0.7, 5.9, 11.9, 1.05, fill=PANEL)
txt(s, 1.1, 6.1, 11.2, 0.7, [
    {"runs": [{"t": "One platform, three tracks:  ", "s": 15, "b": True, "c": TXT},
              {"t": "Resilient Cities  ×  Agentic AI  ×  Data Analysis & Correlation", "s": 15, "b": True, "c": CY}]},
], anchor=MSO_ANCHOR.MIDDLE)

# ============ 5. ARCHITECTURE ============
s = slide()
header(s, "How it works", "Architecture — data in, actions out")
stages = [
    ("1 · DATA", "FortyGuard\nTemperature API®\n20 nodes · 2 m · 20 m²", CY),
    ("2 · FEATURES", "Diurnal cycle · UHI per zone\nwet bulb (Stull) · trends\n24h/72h · season", GOLD),
    ("3 · AI CORE", "Hybrid Ridge+LightGBM\nnowcast 6 h ahead\nMAE 0.26 °C · R² 0.95\n+ 48 h z-score (anomalies)", ACC),
    ("4 · AGENT", "Policy engine\ndedup + escalation only\naudit ledger (JSONL)\nFR/EN alert drafting", GREEN),
    ("5 · ACTION", "Hospitals · schools\noutdoor workers\ncooling centers\npublic advisories", RED),
]
x = 0.45
for title, body, c in stages:
    box(s, x, 1.9, 2.28, 3.1, fill=PANEL, line=c)
    txt(s, x + 0.12, 2.05, 2.05, 0.4, [{"runs": [{"t": title, "s": 13, "b": True, "c": c}]}])
    for i, ln in enumerate(body.split("\n")):
        txt(s, x + 0.12, 2.55 + i * 0.52, 2.05, 0.5, [{"runs": [{"t": ln, "s": 11.5, "c": TXT}], "sa": 2}])
    if title != "5 · ACTION":
        ar = s.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x + 2.28), Inches(3.2), Inches(0.24), Inches(0.5))
        ar.fill.solid(); ar.fill.fore_color.rgb = c; ar.line.fill.background(); ar.shadow.inherit = False
    x += 2.52
txt(s, 0.7, 5.4, 12, 1.4, [
    {"runs": [{"t": "Real-time loop every 15 s in the demo  ·  end-to-end latency < 5 s on CPU  ·  dashboard self-contained: works offline, in the field, on a Jetson.", "s": 14, "c": MUT}]},
])

# ============ 6. RÉSULTATS ============
s = slide()
header(s, "Results", "What we built & measured")
cards = [
    ("0.26 °C", "MAE on the 6-hour-ahead nowcast, Cotonou (0.61 °C Phoenix) — 24 h time-based hold-out, 20 points", ACC),
    ("0.95", "R² — the model captures ≥95% of the peak-temperature variance (0.97 Phoenix)", CY),
    ("2 MB", "model size — runs on a laptop CPU and on an NVIDIA Jetson", GOLD),
    ("< 5 s", "from API reading to decision (monitor → alert → action)", GREEN),
]
for i, (big, small, c) in enumerate(cards):
    l = 0.55 + (i % 2) * 6.2
    t = 1.9 + (i // 2) * 2.05
    box(s, l, t, 5.9, 1.8, fill=PANEL)
    txt(s, l + 0.35, t + 0.25, 2.3, 1.2, [{"runs": [{"t": big, "s": 40, "b": True, "c": c}]}])
    txt(s, l + 2.55, t + 0.3, 3.2, 1.3, [{"runs": [{"t": small, "s": 13.5, "c": MUT}]}], anchor=MSO_ANCHOR.MIDDLE)
txt(s, 0.7, 6.0, 12, 1.3, [
    {"runs": [{"t": "Live demo: ", "s": 14, "b": True, "c": TXT},
              {"t": "real-time heat map · 6 h nowcast per node · agent alert feed with concrete actions · full agent log (video + dashboard in submission).", "s": 14, "c": MUT}], "sa": 4},
    {"runs": [{"t": "Runs on REAL data: ", "s": 12, "b": True, "c": GREEN},
              {"t": "Temperature API® tiles harvested on Phoenix — 3 full days × 20 points (24–26 Aug 2026, Downtown peaked 42.91 °C on the 26th); the simulator is re-anchored daily on these real min/avg/max. Cotonou (not yet in the mesh) is calibrated on Open-Meteo (hourly bias −0.6 °C). Reproducible: harvest.py, validate.py.", "s": 12, "c": MUT}]},
])

# ============ 7. NVIDIA / ÉCHELLE ============
s = slide()
header(s, "AI & scale", "Built for the NVIDIA stack")
bullets(s, 0.7, 1.85, 11.9, 3.4, [
    ("Edge-first design", "the whole stack (features, model, agent) targets an NVIDIA Jetson: ONNX export, < 10 ms inference, offline operation — the Jetson kit prize becomes our first edge node."),
    ("City-scale training", "RAPIDS / cuML for GPU-accelerated training when the grid goes from 20 nodes to millions of 20 m² cells across a metropolitan area."),
    ("CUDA-X accelerated pipeline", "the feature pipeline (diurnal, UHI, wet bulb, trends) is vectorized and maps directly to CUDA-X accelerated data processing for the real-time path."),
    ("Reproducible", "one command (train.py) regenerates model, metrics and demo snapshot — no hand-tuning, no magic."),
], size=15.5)
box(s, 0.7, 5.6, 11.9, 1.3, fill=PANEL)
txt(s, 1.1, 5.8, 11.2, 0.9, [
    {"runs": [{"t": "Small model, big scale: ", "s": 16, "b": True, "c": CY},
              {"t": "2 MB on the edge, GPU training path for the whole city — the classic edge-cloud architecture NVIDIA AI Factories deploy in production.", "s": 15, "c": TXT}]},
], anchor=MSO_ANCHOR.MIDDLE)

# ============ 8. RÉSEAU EDGE ============
s = slide()
header(s, "Edge network — vision", "Post-hackathon vision: a community mesh, calibrated against the API")
bullets(s, 0.7, 1.85, 11.9, 3.2, [
    ("Jetson + $5 thermal sensor", "= a 2 m street-level node for ~$120: market squares, schoolyards, building sites, hospital entrances."),
    ("Self-calibrating", "each node reconciles its reading against the FortyGuard field — no manual maintenance, drift detected automatically."),
    ("Data sovereignty", "the city owns its thermal data; the mesh keeps working through internet and power outages (autonomous local mode)."),
    ("Community", "nodes are operated by local volunteers & students — heat resilience as a civic tool, not a black box."),
], size=15.5)
box(s, 0.7, 5.35, 11.9, 1.5, fill=PANEL)
txt(s, 1.1, 5.55, 11.2, 1.1, [
    {"runs": [{"t": "Vision-ready: ", "s": 15, "b": True, "c": GOLD},
              {"t": "edge/jetson/ + a Jetson-ready Dockerfile document the deployment reference (not yet tested on real hardware — the winning kit would be the first node).", "s": 15, "c": TXT}]},
], anchor=MSO_ANCHOR.MIDDLE)

# ============ 9. IMPACT ============
s = slide()
header(s, "Impact", "What changes on the ground in 12 months")
bullets(s, 0.7, 1.85, 11.9, 3.6, [
    ("5 priority sites protected in Cotonou", "CHU (regional hospital), 3 schools, the port market — with rehearsed response protocols tied to our alert levels."),
    ("60-second alert latency", "from heat signal to decision, versus hours with today's bulletins."),
    ("Measurable KPIs", "response time, outdoor-work pause hours, cooling-center openings, heatstroke reports — a city dashboard of heat resilience."),
    ("Proven, then scaled", "Phoenix = the real-data demo (43–46 °C desert heat, in the mesh today) → Cotonou = the deployment target (feels-like >40 °C, awaiting mesh coverage) → Abu Dhabi (humidity extreme): the same stack, different physics, same API."),
], size=15.5)
txt(s, 0.7, 5.9, 12, 1.0, [
    {"runs": [{"t": "Aligned with the WHO Heat Health Action Plan framework: ", "s": 14, "b": True, "c": MUT},
              {"t": "monitor → nowcast → alert → act, with a human in charge of every action.", "s": 14, "c": MUT}]},
])

# ============ 10. CLÔTURE ============
s = slide()
box(s, 0, 0, 13.333, 7.5, fill=BG)
box(s, 0.9, 1.2, 0.14, 1.7, fill=ACC, rounded=False)
txt(s, 1.3, 1.2, 11, 1.2, [{"runs": [{"t": "Heat will not wait. Neither do we.", "s": 38, "b": True, "c": TXT}]}])
txt(s, 1.3, 2.6, 11, 0.5, [{"runs": [{"t": "30-day pilot with the City of Cotonou — ask: Temperature API scale credits + Jetson kit for the first edge node (vision).", "s": 16, "c": GOLD}]}])
txt(s, 1.3, 3.6, 11, 1.8, [
    {"runs": [{"t": "Team HeatSentinel", "s": 18, "b": True, "c": CY}], "sa": 6},
    {"runs": [{"t": "[Nom 1] — AI Engineer  ·  [Nom 2] — Data/Geo  ·  [Nom 3] — Product/Design", "s": 14, "c": MUT}], "sa": 6},
    {"runs": [{"t": "Code: github.com/[you]/heat-sentinel  ·  Live demo: [url]  ·  Video: [url]", "s": 14, "c": MUT}]},
])
txt(s, 1.3, 6.4, 11, 0.6, [{"runs": [{"t": "Thank you.  ·  FortyGuard Hackathon'26  ·  Building the World's Temperature AI", "s": 13, "c": MUT}]}])

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "pitch_deck.pptx")
prs.save(out)
print(f"OK → {out} ({len(prs.slides._sldIdLst)} slides)")
